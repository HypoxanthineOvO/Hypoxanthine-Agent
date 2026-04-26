from __future__ import annotations

from collections.abc import Iterable
from datetime import UTC, datetime
from pathlib import Path
from typing import Any, Literal

import fitz
from docx import Document
from pptx import Presentation
import structlog
import yaml

from hypo_agent.core.config_loader import get_agent_root, get_memory_dir
from hypo_agent.core.directory_index import (
    build_directory_tree,
    load_directory_index_payload,
    merge_directory_descriptions,
)
from hypo_agent.models import SkillOutput
from hypo_agent.security.permission_manager import PermissionManager
from hypo_agent.skills.base import BaseSkill

logger = structlog.get_logger("hypo_agent.skills.fs_skill")

Operation = Literal["read", "write", "execute"]


class FileSystemSkill(BaseSkill):
    MAX_FILE_CHARS = 16000
    MAX_LIST_ENTRIES = 200

    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".yaml",
        ".yml",
        ".json",
        ".py",
        ".csv",
        ".log",
        ".conf",
        ".ini",
        ".toml",
        ".sh",
    }
    IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp"}

    name = "filesystem"
    description = "Read, write, and inspect files and directories."
    required_permissions = ["filesystem"]

    def __init__(
        self,
        *,
        permission_manager: PermissionManager | None,
        index_file: Path | str | None = None,
    ) -> None:
        self.permission_manager = permission_manager

        if index_file is None:
            index_file = get_memory_dir() / "knowledge" / "directory_index.yaml"

        self.index_file = Path(index_file).expanduser().resolve(strict=False)

    @property
    def tools(self) -> list[dict[str, Any]]:
        return [
            {
                "type": "function",
                "function": {
                    "name": "read_file",
                    "description": (
                        "Read file content with format-aware extraction. Use this first when "
                        "the user asks to inspect or read a file; do not assume access is "
                        "denied before trying. If the path is blocked, the tool will return "
                        "an explicit permission error."
                    ),
                    "parameters": {
                        "type": "object",
                        "properties": {"path": {"type": "string"}},
                        "required": ["path"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "write_file",
                    "description": (
                        "Create or overwrite a text file. Use this for explicit file edits; "
                        "writes require directory write permission and the tool will report "
                        "permission errors explicitly."
                    ),
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string"},
                            "content": {"type": "string"},
                        },
                        "required": ["path", "content"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "list_directory",
                    "description": (
                        "List directory entries with depth support. Use this to inspect a "
                        "directory before claiming you cannot access it; blocked paths will "
                        "return an explicit permission error."
                    ),
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string"},
                            "depth": {"type": "integer", "minimum": 1},
                        },
                        "required": ["path"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "scan_directory",
                    "description": "Scan directory tree and update directory index",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string"},
                            "depth": {"type": "integer", "minimum": 1},
                        },
                        "required": ["path"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "get_directory_index",
                    "description": "Read current directory index YAML",
                    "parameters": {"type": "object", "properties": {}},
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "update_directory_description",
                    "description": "Update description for a directory in index",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string"},
                            "description": {"type": "string"},
                        },
                        "required": ["path", "description"],
                    },
                },
            },
        ]

    async def execute(self, tool_name: str, params: dict[str, Any]) -> SkillOutput:
        if tool_name == "read_file":
            return self._read_file(params)
        if tool_name == "write_file":
            return self._write_file(params)
        if tool_name == "list_directory":
            return self._list_directory(params)
        if tool_name == "scan_directory":
            return self._scan_directory(params)
        if tool_name == "get_directory_index":
            return self._get_directory_index()
        if tool_name == "update_directory_description":
            return self._update_directory_description(params)
        return SkillOutput(
            status="error",
            error_info=f"Unsupported tool '{tool_name}' for filesystem skill",
        )

    def _read_file(self, params: dict[str, Any]) -> SkillOutput:
        raw_path = str(params.get("path", "")).strip()
        if not raw_path:
            return SkillOutput(status="error", error_info="path is required")

        path = self._resolve_read_target(raw_path)
        denied = self._deny_if_no_permission(str(path), "read")
        if denied is not None:
            return denied

        if not path.exists() or not path.is_file():
            return SkillOutput(status="error", error_info=f"File not found: {path}")

        ext = path.suffix.lower()
        try:
            if ext in self.TEXT_EXTENSIONS:
                text = path.read_text(encoding="utf-8", errors="replace")
                content, truncated = self._truncate_text(text)
                size = path.stat().st_size
                logger.info("fs.read", path=str(path), size=size, format="text")
                return SkillOutput(
                    status="success",
                    result=content,
                    metadata={
                        "path": str(path),
                        "size": size,
                        "format": "text",
                        "truncated": truncated,
                    },
                )

            if ext == ".pdf":
                content = self._read_pdf(path)
                text, truncated = self._truncate_text(content)
                size = path.stat().st_size
                logger.info("fs.read", path=str(path), size=size, format="pdf")
                return SkillOutput(
                    status="success",
                    result=text,
                    metadata={
                        "path": str(path),
                        "size": size,
                        "format": "pdf",
                        "truncated": truncated,
                    },
                )

            if ext in self.IMAGE_EXTENSIONS:
                meta = self._image_metadata(path)
                size = path.stat().st_size
                logger.info("fs.read", path=str(path), size=size, format="image")
                return SkillOutput(
                    status="success",
                    result=(
                        f"Image metadata: format={meta['image_format']} "
                        f"width={meta['width']} height={meta['height']} bytes={size}"
                    ),
                    metadata={
                        "path": str(path),
                        "size": size,
                        "format": "image",
                        **meta,
                    },
                )

            if ext == ".pptx":
                content = self._read_pptx(path)
                text, truncated = self._truncate_text(content)
                size = path.stat().st_size
                logger.info("fs.read", path=str(path), size=size, format="pptx")
                return SkillOutput(
                    status="success",
                    result=text,
                    metadata={
                        "path": str(path),
                        "size": size,
                        "format": "pptx",
                        "truncated": truncated,
                    },
                )

            if ext == ".docx":
                content = self._read_docx(path)
                text, truncated = self._truncate_text(content)
                size = path.stat().st_size
                logger.info("fs.read", path=str(path), size=size, format="docx")
                return SkillOutput(
                    status="success",
                    result=text,
                    metadata={
                        "path": str(path),
                        "size": size,
                        "format": "docx",
                        "truncated": truncated,
                    },
                )
        except (OSError, RuntimeError, TypeError, ValueError) as exc:
            message = str(exc)
            if message.startswith("Cannot read encrypted PDF:"):
                return SkillOutput(status="error", error_info=message)
            return SkillOutput(status="error", error_info=f"Failed to read file: {exc}")

        stat = path.stat()
        modified = datetime.fromtimestamp(stat.st_mtime, tz=UTC).isoformat()
        logger.info("fs.read", path=str(path), size=stat.st_size, format="unsupported")
        return SkillOutput(
            status="success",
            result=(
                f"Unsupported format for '{path.name}'. "
                f"size={stat.st_size} modified={modified}"
            ),
            metadata={
                "path": str(path),
                "size": stat.st_size,
                "modified_time": modified,
                "format": "unsupported",
            },
        )

    def _write_file(self, params: dict[str, Any]) -> SkillOutput:
        raw_path = str(params.get("path", "")).strip()
        if not raw_path:
            return SkillOutput(status="error", error_info="path is required")
        content = str(params.get("content", ""))

        denied = self._deny_if_no_permission(raw_path, "write")
        if denied is not None:
            return denied

        path = Path(raw_path).expanduser().resolve(strict=False)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(content, encoding="utf-8")
        size = path.stat().st_size
        logger.info("fs.write", path=str(path), size=size, format="text")
        return SkillOutput(
            status="success",
            result=f"Wrote file: {path}",
            metadata={"path": str(path), "bytes": size},
        )

    def _list_directory(self, params: dict[str, Any]) -> SkillOutput:
        raw_path = str(params.get("path", "")).strip()
        if not raw_path:
            return SkillOutput(status="error", error_info="path is required")
        depth = max(1, int(params.get("depth") or 1))

        denied = self._deny_if_no_permission(raw_path, "read")
        if denied is not None:
            return denied

        root = Path(raw_path).expanduser().resolve(strict=False)
        if not root.exists() or not root.is_dir():
            return SkillOutput(status="error", error_info=f"Directory not found: {root}")

        lines: list[str] = []
        count = 0
        truncated = False
        for item, level in self._iter_entries(root, depth):
            if count >= self.MAX_LIST_ENTRIES:
                truncated = True
                break
            rel = item.relative_to(root) if item != root else Path(".")
            item_type = "dir" if item.is_dir() else "file"
            size = item.stat().st_size if item.is_file() else 0
            lines.append(f"{'  ' * level}{rel}\t{item_type}\t{size}")
            count += 1

        if truncated:
            lines.append(f"... truncated to {self.MAX_LIST_ENTRIES} entries")

        logger.info("fs.list", path=str(root), size=count, format="tree")
        return SkillOutput(
            status="success",
            result="\n".join(lines),
            metadata={
                "path": str(root),
                "depth": depth,
                "count": count,
                "truncated": truncated,
            },
        )

    def _scan_directory(self, params: dict[str, Any]) -> SkillOutput:
        raw_path = str(params.get("path", "")).strip()
        if not raw_path:
            return SkillOutput(status="error", error_info="path is required")
        depth = max(1, int(params.get("depth") or 2))

        denied = self._deny_if_no_permission(raw_path, "read")
        if denied is not None:
            return denied

        root = Path(raw_path).expanduser().resolve(strict=False)
        if not root.exists() or not root.is_dir():
            return SkillOutput(status="error", error_info=f"Directory not found: {root}")

        new_tree = self._scan_tree(root, depth)
        payload = self._load_directory_index()
        directories = payload.setdefault("directories", {})
        if not isinstance(directories, dict):
            directories = {}
            payload["directories"] = directories

        root_key = str(root)
        old_tree = directories.get(root_key, {})
        if isinstance(old_tree, dict):
            new_tree = self._merge_descriptions(new_tree, old_tree)
        directories[root_key] = new_tree
        payload["agent_root"] = str(get_agent_root())
        payload["generated_at"] = datetime.now(UTC).isoformat()
        payload["last_scan"] = datetime.now(UTC).isoformat()

        self.index_file.parent.mkdir(parents=True, exist_ok=True)
        self.index_file.write_text(
            yaml.safe_dump(payload, sort_keys=False, allow_unicode=True),
            encoding="utf-8",
        )

        logger.info(
            "fs.scan",
            path=str(root),
            depth=depth,
            index_path=str(self.index_file),
        )
        return SkillOutput(
            status="success",
            result=f"Scanned directory: {root}",
            metadata={"path": str(root), "depth": depth, "index_path": str(self.index_file)},
        )

    def _get_directory_index(self) -> SkillOutput:
        denied = self._deny_if_no_permission(str(self.index_file), "read")
        if denied is not None:
            return denied

        if not self.index_file.exists():
            return SkillOutput(status="success", result="directories: {}\n")

        content = self.index_file.read_text(encoding="utf-8")
        return SkillOutput(
            status="success",
            result=content,
            metadata={"index_path": str(self.index_file)},
        )

    def _update_directory_description(self, params: dict[str, Any]) -> SkillOutput:
        raw_path = str(params.get("path", "")).strip()
        description = str(params.get("description", ""))
        if not raw_path:
            return SkillOutput(status="error", error_info="path is required")

        denied = self._deny_if_no_permission(raw_path, "read")
        if denied is not None:
            return denied

        if (
            self.permission_manager is not None
            and not self.permission_manager.has_whitelist_match(raw_path)
        ):
            return SkillOutput(
                status="error",
                error_info=(
                    f"Permission denied: Path '{Path(raw_path).expanduser().resolve(strict=False)}' "
                    "is outside explicit whitelist visibility"
                ),
            )

        denied = self._deny_if_no_permission(str(self.index_file), "write")
        if denied is not None:
            return denied

        target = Path(raw_path).expanduser().resolve(strict=False)
        payload = self._load_directory_index()
        directories = payload.get("directories")
        if not isinstance(directories, dict):
            return SkillOutput(status="error", error_info="Directory index is empty")

        node = self._find_node(directories, target)
        if node is None:
            return SkillOutput(
                status="error",
                error_info=f"Directory not found in index: {target}",
            )

        node["description"] = description
        self.index_file.parent.mkdir(parents=True, exist_ok=True)
        self.index_file.write_text(
            yaml.safe_dump(payload, sort_keys=False, allow_unicode=True),
            encoding="utf-8",
        )

        logger.info(
            "fs.index.update",
            path=str(target),
            index_path=str(self.index_file),
        )
        return SkillOutput(
            status="success",
            result=f"Updated description for: {target}",
            metadata={"path": str(target), "description": description},
        )

    def _deny_if_no_permission(
        self,
        path: str,
        operation: Operation,
    ) -> SkillOutput | None:
        if self.permission_manager is None:
            return None

        allowed, reason = self.permission_manager.check_permission(
            path,
            operation,
            log_allowed=False,
        )
        if allowed:
            return None
        return SkillOutput(status="error", error_info=f"Permission denied: {reason}")

    def _resolve_read_target(self, raw_path: str) -> Path:
        requested = Path(raw_path).expanduser()
        candidates: list[Path] = []

        if requested.is_absolute():
            candidates.append(requested.resolve(strict=False))
        else:
            candidates.append(requested.resolve(strict=False))
            memory_dir = get_memory_dir().resolve(strict=False)
            candidates.append((memory_dir / requested).resolve(strict=False))
            if len(requested.parts) == 1:
                candidates.append((memory_dir / "exports" / requested.name).resolve(strict=False))

        seen: set[Path] = set()
        unique_candidates: list[Path] = []
        for candidate in candidates:
            if candidate in seen:
                continue
            seen.add(candidate)
            unique_candidates.append(candidate)

        for candidate in unique_candidates:
            if candidate.exists() and candidate.is_file():
                return candidate
        return unique_candidates[0]

    def _truncate_text(self, content: str) -> tuple[str, bool]:
        if len(content) <= self.MAX_FILE_CHARS:
            return content, False
        return content[: self.MAX_FILE_CHARS] + "\n[truncated]", True

    def _load_directory_index(self) -> dict[str, Any]:
        return load_directory_index_payload(self.index_file)

    def _iter_entries(self, root: Path, max_depth: int) -> Iterable[tuple[Path, int]]:
        queue: list[tuple[Path, int]] = [(root, 0)]
        while queue:
            current, level = queue.pop(0)
            if current != root:
                yield current, level - 1

            if level >= max_depth or not current.is_dir():
                continue

            children = sorted(current.iterdir(), key=lambda item: (item.is_file(), item.name))
            for child in children:
                queue.append((child, level + 1))

    def _read_pdf(self, path: Path) -> str:
        with fitz.open(path) as document:
            if document.is_encrypted:
                raise ValueError(f"Cannot read encrypted PDF: {path}")
            texts = [page.get_text("text") for page in document]
        return "\n".join(texts).strip()

    def _read_docx(self, path: Path) -> str:
        document = Document(path)
        lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n".join(lines).strip()

    def _read_pptx(self, path: Path) -> str:
        presentation = Presentation(path)
        lines: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    slide_lines.append(text.strip())
            if slide_lines:
                lines.append(f"Slide {index}:")
                lines.extend(slide_lines)
        return "\n".join(lines).strip()

    def _image_metadata(self, path: Path) -> dict[str, Any]:
        pixmap = fitz.Pixmap(str(path))
        try:
            return {
                "image_format": path.suffix.lower().lstrip("."),
                "width": pixmap.width,
                "height": pixmap.height,
            }
        finally:
            del pixmap

    def _scan_tree(self, root: Path, depth: int) -> dict[str, Any]:
        return build_directory_tree(root, agent_root=get_agent_root(), max_depth=depth)

    def _merge_descriptions(
        self,
        fresh: dict[str, Any],
        existing: dict[str, Any],
    ) -> dict[str, Any]:
        return merge_directory_descriptions(fresh, existing)

    def _find_node(
        self,
        directories: dict[str, Any],
        target_path: Path,
    ) -> dict[str, Any] | None:
        for root_key, root_node in directories.items():
            if not isinstance(root_node, dict):
                continue
            root_path = Path(root_key).expanduser().resolve(strict=False)
            if target_path == root_path:
                return root_node
            try:
                relative = target_path.relative_to(root_path)
            except ValueError:
                continue

            node = root_node
            for part in relative.parts:
                children = node.get("children")
                if not isinstance(children, dict) or part not in children:
                    node = None
                    break
                child_node = children[part]
                if not isinstance(child_node, dict):
                    node = None
                    break
                node = child_node
            if isinstance(node, dict):
                return node
        return None
