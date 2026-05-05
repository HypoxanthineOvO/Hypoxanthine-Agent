from __future__ import annotations

import asyncio
import base64
from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation
import yaml

from hypo_agent.models import DirectoryWhitelist, WhitelistRule
from hypo_agent.security.permission_manager import PermissionManager
from hypo_agent.skills import fs_skill as fs_module
from hypo_agent.skills.fs_skill import FileSystemSkill


def _build_skill(tmp_path: Path) -> tuple[FileSystemSkill, Path, Path]:
    writable = tmp_path / "workspace"
    readonly = tmp_path / "readonly"
    index_file = tmp_path / "memory" / "knowledge" / "directory_index.yaml"
    writable.mkdir(parents=True, exist_ok=True)
    readonly.mkdir(parents=True, exist_ok=True)
    index_file.parent.mkdir(parents=True, exist_ok=True)

    manager = PermissionManager(
        DirectoryWhitelist(
            rules=[
                WhitelistRule(path=str(writable), permissions=["read", "write"]),
                WhitelistRule(path=str(readonly), permissions=["read"]),
                WhitelistRule(
                    path=str(index_file.parent),
                    permissions=["read", "write"],
                ),
            ],
            default_policy="readonly",
        )
    )
    skill = FileSystemSkill(
        permission_manager=manager,
        index_file=index_file,
    )
    return skill, writable, readonly


def test_fs_skill_declares_required_permissions(tmp_path: Path) -> None:
    skill, _, _ = _build_skill(tmp_path)
    assert skill.required_permissions == ["filesystem"]


def test_fs_skill_tool_descriptions_tell_model_to_try_reads_first(tmp_path: Path) -> None:
    skill, _, _ = _build_skill(tmp_path)
    descriptions = {
        tool["function"]["name"]: tool["function"]["description"]
        for tool in skill.tools
    }

    assert "do not assume access is denied before trying" in descriptions["read_file"]
    assert "permission error" in descriptions["read_file"]
    assert "writes require directory write permission" in descriptions["write_file"]
    assert "before claiming you cannot access it" in descriptions["list_directory"]


def test_blocked_path_filesystem_denied(tmp_path: Path) -> None:
    blocked = tmp_path / "blocked.txt"
    blocked.write_text("secret", encoding="utf-8")
    manager = PermissionManager(
        DirectoryWhitelist(
            rules=[],
            default_policy="readonly",
            blocked_paths=[str(blocked)],
        )
    )
    skill = FileSystemSkill(permission_manager=manager, index_file=tmp_path / "index.yaml")

    output = asyncio.run(skill.execute("read_file", {"path": str(blocked)}))

    assert output.status == "error"
    assert "permission" in (output.error_info or "").lower()


def test_gray_zone_read_allowed(tmp_path: Path) -> None:
    gray = tmp_path / "gray.txt"
    gray.write_text("ok", encoding="utf-8")
    manager = PermissionManager(DirectoryWhitelist(rules=[], default_policy="readonly"))
    skill = FileSystemSkill(permission_manager=manager, index_file=tmp_path / "index.yaml")

    output = asyncio.run(skill.execute("read_file", {"path": str(gray)}))

    assert output.status == "success"


def test_gray_zone_write_denied(tmp_path: Path) -> None:
    gray = tmp_path / "gray.txt"
    manager = PermissionManager(DirectoryWhitelist(rules=[], default_policy="readonly"))
    skill = FileSystemSkill(permission_manager=manager, index_file=tmp_path / "index.yaml")

    output = asyncio.run(skill.execute("write_file", {"path": str(gray), "content": "x"}))

    assert output.status == "error"


def test_read_file_truncates_large_text(tmp_path: Path) -> None:
    skill, _, readonly = _build_skill(tmp_path)
    file_path = readonly / "big.log"
    file_path.write_text("a" * 17000, encoding="utf-8")

    output = asyncio.run(skill.execute("read_file", {"path": str(file_path)}))

    assert output.status == "success"
    assert isinstance(output.result, str)
    assert len(output.result) <= FileSystemSkill.MAX_FILE_CHARS + 64
    assert output.metadata["format"] == "text"
    assert output.metadata["truncated"] is True


def test_read_file_extracts_pdf_text(tmp_path: Path) -> None:
    skill, _, readonly = _build_skill(tmp_path)
    file_path = readonly / "sample.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "hello pdf")
    pdf.save(file_path)
    pdf.close()

    output = asyncio.run(skill.execute("read_file", {"path": str(file_path)}))

    assert output.status == "success"
    assert "hello pdf" in output.result.lower()
    assert output.metadata["format"] == "pdf"


def test_read_file_extracts_docx_text(tmp_path: Path) -> None:
    skill, _, readonly = _build_skill(tmp_path)
    file_path = readonly / "sample.docx"
    doc = Document()
    doc.add_paragraph("hello docx")
    doc.save(file_path)

    output = asyncio.run(skill.execute("read_file", {"path": str(file_path)}))

    assert output.status == "success"
    assert "hello docx" in output.result.lower()
    assert output.metadata["format"] == "docx"


def test_read_file_extracts_pptx_text(tmp_path: Path) -> None:
    skill, _, readonly = _build_skill(tmp_path)
    file_path = readonly / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "slide title"
    slide.placeholders[1].text = "slide body"
    presentation.save(file_path)

    output = asyncio.run(skill.execute("read_file", {"path": str(file_path)}))

    assert output.status == "success"
    assert "slide title" in output.result.lower()
    assert "slide body" in output.result.lower()
    assert output.metadata["format"] == "pptx"


def test_read_file_returns_image_metadata_without_ocr(tmp_path: Path) -> None:
    skill, _, readonly = _build_skill(tmp_path)
    file_path = readonly / "pixel.png"
    png_bytes = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+tmR0AAAAASUVORK5CYII="
    )
    file_path.write_bytes(png_bytes)

    output = asyncio.run(skill.execute("read_file", {"path": str(file_path)}))

    assert output.status == "success"
    assert output.metadata["format"] == "image"
    assert output.metadata["image_format"] == "png"
    assert output.metadata["width"] == 1
    assert output.metadata["height"] == 1


def test_read_file_falls_back_for_unsupported_format(tmp_path: Path) -> None:
    skill, _, readonly = _build_skill(tmp_path)
    file_path = readonly / "blob.bin"
    file_path.write_bytes(b"\x00\x01\x02")

    output = asyncio.run(skill.execute("read_file", {"path": str(file_path)}))

    assert output.status == "success"
    assert "unsupported" in output.result.lower()
    assert output.metadata["format"] == "unsupported"


def test_read_file_resolves_exported_attachment_filename_from_memory_exports(
    tmp_path: Path,
    monkeypatch,
) -> None:
    skill, _, _ = _build_skill(tmp_path)
    exported = tmp_path / "memory" / "exports" / "tool-output.md"
    exported.parent.mkdir(parents=True, exist_ok=True)
    exported.write_text("export body", encoding="utf-8")
    monkeypatch.setattr(fs_module, "get_memory_dir", lambda: tmp_path / "memory")

    output = asyncio.run(skill.execute("read_file", {"path": "tool-output.md"}))

    assert output.status == "success"
    assert output.result == "export body"
    assert output.metadata["path"] == str(exported)


def test_read_file_returns_resource_recovery_action_for_missing_file(tmp_path: Path) -> None:
    skill, writable, _ = _build_skill(tmp_path)
    nearby = writable / "channel-report.md"
    older = writable / "old-report.md"
    nearby.write_text("report", encoding="utf-8")
    older.write_text("old", encoding="utf-8")

    output = asyncio.run(skill.execute("read_file", {"path": "report.md"}))

    assert output.status == "error"
    assert "File not found" in output.error_info
    assert output.metadata["resource_resolution"]["status"] == "ambiguous"
    assert output.metadata["recovery_action"]["type"] == "ask_user"
    assert output.metadata["recovery_action"]["reason"] == "multiple_candidates"
    assert output.metadata["resource_candidates"][0]["display_name"] == "channel-report.md"


def test_write_file_allows_whitelisted_path(tmp_path: Path) -> None:
    skill, writable, _ = _build_skill(tmp_path)
    file_path = writable / "nested" / "notes.md"

    output = asyncio.run(
        skill.execute(
            "write_file",
            {"path": str(file_path), "content": "hello"},
        )
    )

    assert output.status == "success"
    assert file_path.read_text(encoding="utf-8") == "hello"
    assert output.metadata["bytes"] == 5


def test_write_file_denies_outside_whitelist(tmp_path: Path) -> None:
    skill, _, _ = _build_skill(tmp_path)
    file_path = tmp_path / "outside" / "notes.md"

    output = asyncio.run(
        skill.execute(
            "write_file",
            {"path": str(file_path), "content": "hello"},
        )
    )

    assert output.status == "error"
    assert "permission denied" in output.error_info.lower()


def test_write_file_denies_path_traversal(tmp_path: Path) -> None:
    skill, writable, _ = _build_skill(tmp_path)
    file_path = writable / ".." / "escape.txt"

    output = asyncio.run(
        skill.execute(
            "write_file",
            {"path": str(file_path), "content": "escape"},
        )
    )

    assert output.status == "error"
    assert "outside whitelist" in output.error_info.lower()


def test_list_directory_caps_result_entries(tmp_path: Path) -> None:
    skill, writable, _ = _build_skill(tmp_path)
    for idx in range(210):
        (writable / f"file-{idx}.txt").write_text(str(idx), encoding="utf-8")

    output = asyncio.run(skill.execute("list_directory", {"path": str(writable), "depth": 1}))

    assert output.status == "success"
    assert output.metadata["count"] == 200
    assert output.metadata["truncated"] is True


def test_scan_directory_writes_index_file(tmp_path: Path) -> None:
    skill, writable, _ = _build_skill(tmp_path)
    (writable / "alpha").mkdir(parents=True, exist_ok=True)
    (writable / "alpha" / "a.txt").write_text("a", encoding="utf-8")
    (writable / "b.txt").write_text("b", encoding="utf-8")

    output = asyncio.run(skill.execute("scan_directory", {"path": str(writable), "depth": 2}))

    assert output.status == "success"
    assert skill.index_file.exists()
    payload = yaml.safe_load(skill.index_file.read_text(encoding="utf-8"))
    root_key = str(writable.resolve(strict=False))
    assert root_key in payload["directories"]
    assert payload["directories"][root_key]["children"]["alpha"]["file_count"] == 1
    assert "last_scan" in payload


def test_get_directory_index_returns_content(tmp_path: Path) -> None:
    skill, _, _ = _build_skill(tmp_path)
    skill.index_file.parent.mkdir(parents=True, exist_ok=True)
    skill.index_file.write_text(
        "directories:\n  /tmp/example:\n    description: \"\"\n    file_count: 0\n    children: {}\n",
        encoding="utf-8",
    )

    output = asyncio.run(skill.execute("get_directory_index", {}))

    assert output.status == "success"
    assert "/tmp/example" in output.result


def test_update_directory_description_updates_index(tmp_path: Path) -> None:
    skill, writable, _ = _build_skill(tmp_path)
    (writable / "a.txt").write_text("a", encoding="utf-8")
    asyncio.run(skill.execute("scan_directory", {"path": str(writable), "depth": 2}))

    output = asyncio.run(
        skill.execute(
            "update_directory_description",
            {"path": str(writable), "description": "workspace folder"},
        )
    )

    assert output.status == "success"
    payload = yaml.safe_load(skill.index_file.read_text(encoding="utf-8"))
    root_key = str(writable.resolve(strict=False))
    assert payload["directories"][root_key]["description"] == "workspace folder"


def test_scan_directory_preserves_existing_description(tmp_path: Path) -> None:
    skill, writable, _ = _build_skill(tmp_path)
    (writable / "a.txt").write_text("a", encoding="utf-8")
    asyncio.run(skill.execute("scan_directory", {"path": str(writable), "depth": 2}))
    asyncio.run(
        skill.execute(
            "update_directory_description",
            {"path": str(writable), "description": "keep me"},
        )
    )

    (writable / "b.txt").write_text("b", encoding="utf-8")
    output = asyncio.run(skill.execute("scan_directory", {"path": str(writable), "depth": 2}))

    assert output.status == "success"
    payload = yaml.safe_load(skill.index_file.read_text(encoding="utf-8"))
    root_key = str(writable.resolve(strict=False))
    assert payload["directories"][root_key]["description"] == "keep me"


def test_fs_skill_emits_observability_events(tmp_path: Path, monkeypatch) -> None:
    skill, writable, readonly = _build_skill(tmp_path)
    (readonly / "note.txt").write_text("hello", encoding="utf-8")
    (writable / "scan-me.txt").write_text("x", encoding="utf-8")
    events: list[str] = []

    class LogRecorder:
        def info(self, event: str, **kwargs) -> None:
            del kwargs
            events.append(event)

        def warning(self, event: str, **kwargs) -> None:
            del kwargs
            events.append(event)

    monkeypatch.setattr(fs_module, "logger", LogRecorder())

    asyncio.run(skill.execute("read_file", {"path": str(readonly / "note.txt")}))
    asyncio.run(skill.execute("write_file", {"path": str(writable / "w.txt"), "content": "c"}))
    asyncio.run(skill.execute("list_directory", {"path": str(writable), "depth": 1}))
    asyncio.run(skill.execute("scan_directory", {"path": str(writable), "depth": 2}))
    asyncio.run(
        skill.execute(
            "update_directory_description",
            {"path": str(writable), "description": "desc"},
        )
    )

    assert "fs.read" in events
    assert "fs.write" in events
    assert "fs.list" in events
    assert "fs.scan" in events
    assert "fs.index.update" in events


def test_update_description_allows_readonly_target_if_index_file_writable(tmp_path: Path) -> None:
    src_root = tmp_path / "src"
    target_dir = src_root / "hypo_agent" / "security"
    target_dir.mkdir(parents=True, exist_ok=True)
    (target_dir / "guard.py").write_text("x = 1\n", encoding="utf-8")
    index_file = tmp_path / "memory" / "knowledge" / "directory_index.yaml"

    manager = PermissionManager(
        DirectoryWhitelist(
            rules=[
                WhitelistRule(path=str(src_root), permissions=["read"]),
                WhitelistRule(
                    path=str(index_file.parent),
                    permissions=["read", "write"],
                ),
            ],
            default_policy="readonly",
        )
    )
    skill = FileSystemSkill(permission_manager=manager, index_file=index_file)

    asyncio.run(skill.execute("scan_directory", {"path": str(src_root), "depth": 4}))
    output = asyncio.run(
        skill.execute(
            "update_directory_description",
            {"path": str(target_dir), "description": "security module"},
        )
    )

    assert output.status == "success"
    payload = yaml.safe_load(index_file.read_text(encoding="utf-8"))
    root_key = str(src_root.resolve(strict=False))
    assert payload["directories"][root_key]["children"]["hypo_agent"]["children"][
        "security"
    ]["description"] == "security module"


def test_update_description_rejects_path_outside_explicit_whitelist(tmp_path: Path) -> None:
    src_root = tmp_path / "src"
    src_root.mkdir(parents=True, exist_ok=True)
    index_file = tmp_path / "memory" / "knowledge" / "directory_index.yaml"
    outside = tmp_path / "outside" / "secret"
    outside.mkdir(parents=True, exist_ok=True)

    manager = PermissionManager(
        DirectoryWhitelist(
            rules=[
                WhitelistRule(path=str(src_root), permissions=["read"]),
                WhitelistRule(
                    path=str(index_file.parent),
                    permissions=["read", "write"],
                ),
            ],
            default_policy="readonly",
        )
    )
    skill = FileSystemSkill(permission_manager=manager, index_file=index_file)
    asyncio.run(skill.execute("scan_directory", {"path": str(src_root), "depth": 2}))

    output = asyncio.run(
        skill.execute(
            "update_directory_description",
            {"path": str(outside), "description": "should fail"},
        )
    )

    assert output.status == "error"
    assert "explicit whitelist" in output.error_info.lower()


def test_list_directory_handles_mixed_file_and_directory_entries(tmp_path: Path) -> None:
    skill, writable, _ = _build_skill(tmp_path)
    (writable / "top.txt").write_text("top", encoding="utf-8")
    nested = writable / "nested"
    nested.mkdir(parents=True, exist_ok=True)
    (nested / "child.txt").write_text("child", encoding="utf-8")

    output = asyncio.run(skill.execute("list_directory", {"path": str(writable), "depth": 2}))

    assert output.status == "success"
    assert "top.txt" in output.result
    assert "nested" in output.result
    assert "child.txt" in output.result


def test_read_file_returns_friendly_error_for_encrypted_pdf(tmp_path: Path) -> None:
    skill, _, readonly = _build_skill(tmp_path)
    file_path = readonly / "encrypted.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "secret")
    pdf.save(
        file_path,
        encryption=fitz.PDF_ENCRYPT_AES_256,
        owner_pw="owner-password",
        user_pw="user-password",
    )
    pdf.close()

    output = asyncio.run(skill.execute("read_file", {"path": str(file_path)}))

    assert output.status == "error"
    assert f"Cannot read encrypted PDF: {file_path}" in output.error_info


def test_fs_skill_allows_agent_memory_and_config_and_home_read_but_denies_src_write(
    tmp_path: Path,
) -> None:
    repo_root = tmp_path / "Hypo-Agent"
    memory_dir = repo_root / "memory"
    config_dir = repo_root / "config"
    src_dir = repo_root / "src"
    memory_dir.mkdir(parents=True, exist_ok=True)
    config_dir.mkdir(parents=True, exist_ok=True)
    src_dir.mkdir(parents=True, exist_ok=True)

    manager = PermissionManager(
        DirectoryWhitelist(
            rules=[
                WhitelistRule(path=str(repo_root), permissions=["read"]),
                WhitelistRule(path="/home/heyx", permissions=["read"]),
                WhitelistRule(path=str(config_dir), permissions=["read", "write"]),
                WhitelistRule(path=str(memory_dir), permissions=["read", "write"]),
            ],
            default_policy="readonly",
        )
    )
    skill = FileSystemSkill(
        permission_manager=manager,
        index_file=memory_dir / "knowledge" / "directory_index.yaml",
    )

    home_file = tmp_path / "homefile.txt"
    home_file.write_text("hello from home", encoding="utf-8")
    home_target = Path("/home/heyx/test-readme.txt")

    read_home = asyncio.run(skill.execute("read_file", {"path": str(home_file)}))
    write_memory = asyncio.run(
        skill.execute(
            "write_file",
            {"path": str(memory_dir / "notes.md"), "content": "memory ok"},
        )
    )
    write_config = asyncio.run(
        skill.execute(
            "write_file",
            {"path": str(config_dir / "persona.yaml"), "content": "name: Hypo\n"},
        )
    )
    deny_src = asyncio.run(
        skill.execute(
            "write_file",
            {"path": str(src_dir / "app.py"), "content": "print('nope')\n"},
        )
    )

    allowed_home, _ = manager.check_permission(str(home_target), "read")

    assert read_home.status == "success"
    assert read_home.result == "hello from home"
    assert write_memory.status == "success"
    assert (memory_dir / "notes.md").read_text(encoding="utf-8") == "memory ok"
    assert write_config.status == "success"
    assert (config_dir / "persona.yaml").read_text(encoding="utf-8") == "name: Hypo\n"
    assert allowed_home is True
    assert deny_src.status == "error"
    assert "permission denied" in (deny_src.error_info or "").lower()
